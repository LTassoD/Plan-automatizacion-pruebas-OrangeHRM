from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BLUE = RGBColor(0x1A, 0x3C, 0x6E)
ORANGE = RGBColor(0xFF, 0x7B, 0x1D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
RED = RGBColor(0xF4, 0x43, 0x36)


def add_bg(slide, color=DARK_BLUE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_box(slide, text, left=0.5, top=0.3, width=12, height=1, font_size=36, color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = True
    return txBox


def add_text_box(slide, text, left=0.5, top=1.5, width=12, height=5, font_size=18, color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return txBox


# ─── Slide 1: Portada ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_title_box(slide, "EVALUACIÓN PARCIAL 3", 1, 0.5, 11, 1.5, 44, ORANGE)
add_text_box(slide,
    "Analizando los Resultados Obtenidos\n\n"
    "Automatización de Pruebas - OrangeHRM\n\n"
    "Estudiante: Luis Tasso\n"
    "Sección: 802V\n"
    "Profesor: Manuel Soto\n\n"
    "API0101 - Automatización de Pruebas",
    1, 2.5, 11, 4.5, 22, WHITE)

# ─── Slide 2: Agenda ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_box(slide, "Agenda", 0.5, 0.3, 12, 1, 36, DARK_BLUE)
add_text_box(slide,
    "1. Ejecución de Pruebas y Registro de Evidencias\n"
    "2. Evaluación de Resultados\n"
    "3. Análisis de Incidencias y Oportunidades de Mejora\n"
    "4. Métricas de Calidad\n"
    "5. Propuesta de Mejora\n"
    "6. Conclusiones",
    0.5, 1.5, 12, 5, 24, DARK_BLUE)

# ─── Slide 3: Resumen del Proyecto ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_box(slide, "Resumen del Proyecto", 0.5, 0.3, 12, 1, 36, DARK_BLUE)
add_text_box(slide,
    "• Sitio: opensource-demo.orangehrmlive.com\n"
    "• Framework: Behave + Selenium WebDriver\n"
    "• Lenguaje: Python 3.14\n"
    "• Data-Driven: Excel (openpyxl)\n"
    "• Navegador: Google Chrome (headless/visual)\n"
    "• 10 archivos .feature | 29 escenarios | 15 casos de prueba",
    0.5, 1.5, 12, 5, 22, DARK_BLUE)

# ─── Slide 4: Ejecución de Pruebas ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_title_box(slide, "Ejecución de Pruebas", 0.5, 0.3, 12, 1, 36, ORANGE)
text = (
    "Módulos cubiertos:\n"
    "  • Login y autenticación (Casos 1, 2, 4, 10)\n"
    "  • Gestión de sesión / Logout (Caso 3)\n"
    "  • Navegación entre módulos (Casos 7, 8, 9)\n"
    "  • PIM - CRUD de empleados (Casos 5-8, 11)\n"
    "  • Perfil My Info (Caso 13)\n"
    "  • Solicitud de licencias (Caso 14)\n\n"
    "Evidencias: 133 screenshots en carpeta evidencias/"
)
add_text_box(slide, text, 0.5, 1.5, 12, 5.5, 20, WHITE)

# ─── Slide 5: Resultados ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_box(slide, "Resultados de Ejecución", 0.5, 0.3, 12, 1, 36, DARK_BLUE)
text = (
    "Total escenarios: 29\n\n"
    "Aprobados: 28 (96.55%)\n\n"
    "Fallidos: 1 (3.45%) - Intencional (Caso 4: assert de título)\n\n"
    "Features ejecutadas: 10 de 10\n\n"
    "Tiempo total: ~7 minutos (suite completa)"
)
add_text_box(slide, text, 0.5, 1.5, 12, 5, 24, DARK_BLUE)

# ─── Slide 6: Métricas ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_title_box(slide, "Métricas de Calidad (PPT 3.3.1)", 0.5, 0.3, 12, 1, 36, ORANGE)
text = (
    "% Automatizable: 100% (29/29 casos)\n"
    "Progreso de Automatización: 100% (29/29 escenarios)\n"
    "Productividad de Diseño: 7.25 escenarios/hora\n"
    "Productividad de Ejecución: 0.067 escenarios/segundo\n"
    "Tasa de Fallos: 3.45% (1 fallo intencional)"
)
add_text_box(slide, text, 0.5, 1.5, 12, 5, 22, WHITE)

# ─── Slide 7: Incidencias Detectadas ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_box(slide, "Incidencias Detectadas", 0.5, 0.3, 12, 1, 36, DARK_BLUE)
text = (
    "• INC-01: Dashboard en chino (\"仪表盘\" en vez de \"Dashboard\")\n"
    "• INC-02: Menú lateral en chino (Leave→休假, Admin→管理员)\n"
    "• INC-03: Timeout en secciones PIM (>20s)\n"
    "• INC-04: Botón Save interceptado por spinner\n"
    "• INC-05: Título de página incorrecto (Caso 4 intencional)\n\n"
    "Causa raíz: Configuración regional del servidor demo + Framework OX"
)
add_text_box(slide, text, 0.5, 1.5, 12, 5.5, 20, DARK_BLUE)

# ─── Slide 8: Oportunidades de Mejora ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_title_box(slide, "Oportunidades de Mejora", 0.5, 0.3, 12, 1, 36, ORANGE)
text = (
    "1. Internacionalización: Detector automático de idioma\n"
    "2. Estabilidad: Localizadores robustos (CSS + fallback JS)\n"
    "3. Esperas: time.sleep() + WebDriverWait combinados\n"
    "4. Page Object Model: Separar localización de negocio\n"
    "5. CI/CD: GitHub Actions para ejecución automática\n"
    "6. Reportes: Allure Reports para dashboards visuales"
)
add_text_box(slide, text, 0.5, 1.5, 12, 5, 22, WHITE)

# ─── Slide 9: Propuestas de Mejora ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_box(slide, "Propuestas de Mejora", 0.5, 0.3, 12, 1, 36, DARK_BLUE)
text = (
    "• Alta prioridad: Detector de idioma automático\n"
    "• Media prioridad: Ejecución paralela (behave -j 4)\n"
    "• Media prioridad: Versionar Excel seed data\n"
    "• Media prioridad: Centralizar lógica de spinner\n"
    "• Baja prioridad: Actualizar valor esperado del título\n"
    "• Técnica: Implementar Page Object Model completo"
)
add_text_box(slide, text, 0.5, 1.5, 12, 5, 22, DARK_BLUE)

# ─── Slide 10: Lecciones Aprendidas ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_title_box(slide, "Lecciones Aprendidas", 0.5, 0.3, 12, 1, 36, ORANGE)
text = (
    "• Localizadores robustos = claves del éxito en automatización\n"
    "• Data-Driven Testing permite probar + combinaciones con - código\n"
    "• Screenshots automáticos = trazabilidad completa\n"
    "• La internacionalización es un desafío real en QA\n"
    "• WebDriverWait > time.sleep() para estabilidad\n"
    "• La suite está lista para CI/CD"
)
add_text_box(slide, text, 0.5, 1.5, 12, 5, 20, WHITE)

# ─── Slide 11: Conclusiones ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_box(slide, "Conclusiones", 0.5, 0.3, 12, 1, 36, DARK_BLUE)
text = (
    "• 29/29 escenarios ejecutados exitosamente\n"
    "• 96.55% de aprobación (fallo intencional documentado)\n"
    "• Suite robusta: soporta inglés y chino\n"
    "• Data-Driven funcional con Excel\n"
    "• 133 evidencias visuales generadas\n"
    "• Reporte HTML con métricas de calidad\n\n"
    "La automatización de pruebas para OrangeHRM\n"
    "cumple con los objetivos y está lista para producción."
)
add_text_box(slide, text, 0.5, 1.5, 12, 5.5, 20, DARK_BLUE)

# ─── Slide 12: Gracias ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_title_box(slide, "Gracias", 1, 2, 11, 1.5, 48, ORANGE)
add_text_box(slide,
    "Luis Tasso | Sección 802V\n"
    "API0101 - Automatización de Pruebas\n"
    "Prof. Manuel Soto",
    1, 4, 11, 2, 24, WHITE)

prs.save("PRESENTACION_EVALUACION_3.pptx")
print("PPT generado: PRESENTACION_EVALUACION_3.pptx")
