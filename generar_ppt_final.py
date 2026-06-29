#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Genera PRESENTACION_FINAL_EFT.pptx para Evaluacion Final Transversal"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT = os.path.join(BASE_DIR, "PRESENTACION_FINAL_EFT.pptx")
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(0x2C, 0x3E, 0x50)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xE7, 0x4C, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x7F, 0x8C, 0x8D)
LIGHT = RGBColor(0xEC, 0xF0, 0xF1)
DARK = RGBColor(0x2C, 0x3E, 0x50)
ACCENT = RGBColor(0x29, 0x80, 0xB9)
GREEN_BG = RGBColor(0xE8, 0xF8, 0xF5)
RED_BG = RGBColor(0xFD, 0xED, 0xEC)
GRAY_TXT = RGBColor(0xBD, 0xC3, 0xC7)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)

def bg(sl, c=WHITE):
    f = sl.background.fill
    f.solid()
    f.fore_color.rgb = c

def rect(sl, l, t, w, h, c, text="", fs=14, b=False, fc=WHITE, a=PP_ALIGN.LEFT):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = c
    s.line.fill.background()
    if text:
        tf = s.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(fs)
        p.font.color.rgb = fc
        p.font.bold = b
        p.alignment = a
    return s

def tb(sl, l, t, w, h, text, fs=18, b=False, c=BLUE, a=PP_ALIGN.LEFT):
    s = sl.shapes.add_textbox(l, t, w, h)
    tf = s.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(fs)
    p.font.color.rgb = c
    p.font.bold = b
    p.alignment = a
    return s

def bullets(sl, l, t, w, h, items, fs=14, c=BLUE):
    s = sl.shapes.add_textbox(l, t, w, h)
    tf = s.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(fs)
        p.font.color.rgb = c
        p.space_after = Pt(6)
    return s

def tbl(sl, l, t, w, h, headers, rows, cw=None):
    ts = sl.shapes.add_table(len(rows)+1, len(headers), l, t, w, h)
    t = ts.table
    if cw:
        for i, w_ in enumerate(cw):
            t.columns[i].width = w_
    for i, h_ in enumerate(headers):
        c = t.cell(0, i)
        c.text = h_
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
        c.fill.solid()
        c.fill.fore_color.rgb = BLUE
    for ri, rd in enumerate(rows):
        for ci, v in enumerate(rd):
            c = t.cell(ri+1, ci)
            c.text = str(v)
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = BLUE
            if ri % 2 == 0:
                c.fill.solid()
                c.fill.fore_color.rgb = LIGHT
    return ts

# ========== SLIDE 1: PORTADA ==========
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, DARK)
rect(sl, Inches(0), Inches(2.5), Inches(13.333), Inches(0.08), ACCENT)
tb(sl, Inches(1), Inches(0.6), Inches(11), Inches(1), "EVALUACION FINAL TRANSVERSAL", 40, True, WHITE, PP_ALIGN.CENTER)
tb(sl, Inches(1), Inches(1.5), Inches(11), Inches(0.8), "Automatizacion de Pruebas - OrangeHRM", 28, False, GRAY_TXT, PP_ALIGN.CENTER)
tb(sl, Inches(1), Inches(3.2), Inches(11), Inches(3),
   "Behave 1.2.6  +  Selenium 4.28  +  Python 3.14\n\n"
   "Estudiante: Luis Tasso  |  Seccion: 802V\n"
   "Profesor: Manuel Soto\n"
   "API0101 - Automatizacion de Pruebas",
   18, False, GRAY_TXT, PP_ALIGN.CENTER)

# ========== SLIDE 2: AGENDA ==========
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE)
tb(sl, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8), "Agenda", 32, True, WHITE)
bullets(sl, Inches(0.8), Inches(1.6), Inches(11), Inches(5), [
    "1. Plan de Pruebas y Cronograma por Sprint",
    "2. Trazabilidad Requisito -> Escenario -> Resultado",
    "3. Resultados de Ejecucion - 29 escenarios, evidencias, fallo intencional",
    "4. Metricas de Calidad (PPT 3.3.1) y Rendimiento",
    "5. Oportunidades de Mejora (A+B+C)",
    "6. Propuestas: Producto (OrangeHRM) vs. Proceso (Framework)",
    "7. Stack Tecnologico: Behave + Python + Selenium",
    "8. Conclusiones y Proyecciones",
], 18)

# ========== SLIDE 3: PLAN DE PRUEBAS Y CRONOGRAMA ==========
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE)
tb(sl, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
   "Plan de Pruebas y Cronograma por Sprint", 28, True, WHITE)

tbl(sl, Inches(0.4), Inches(1.5), Inches(12.5), Inches(2.5),
    ["Sprint", "Semanas", "Features", "Escenarios", "Entregable"],
    [["Sprint 1\nFundamentos", "1-4", "login.feature\nsesion.feature", 4, "Steps genericos\nenvironment.py\nScreenshots"],
     ["Sprint 2\nPIM CRUD + DD", "5-8", "empleados.feature\nbusqueda.feature\nedicion.feature\neliminacion.feature", 15, "ExcelUtils\nData-Driven Steps\n6 Excel testData"],
     ["Sprint 3\nPerfil + Licencias", "9-12", "perfil.feature\nlicencias.feature\nnavigation.feature", 10, "Reportes HTML/JSON\nMetricas\nInforme Final"]],
    cw=[Inches(1.8), Inches(1.5), Inches(3.5), Inches(1.5), Inches(4.2)])

tb(sl, Inches(0.8), Inches(4.3), Inches(11), Inches(2.5),
   "Caso: OrangeHRM (opensource-demo.orangehrmlive.com)\n"
   "Modulos evaluados: Login, Dashboard, PIM, Leave, Admin, My Info, Sesion\n"
   "Tecnicas: BDD (Gherkin), Data-Driven Testing con Excel, Captura automatica de evidencia\n"
   "Tipos de pruebas: Funcionales, Regresion, Humo, Internacionalizacion, Validacion\n\n"
   "Total: 29 escenarios | 10 features | 6 archivos Excel | 177 screenshots",
   14, False, BLUE)

# ========== SLIDE 4: TRAZABILIDAD Y CODIFICACION ==========
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE)
tb(sl, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
   "Trazabilidad: Requisito -> Funcionalidad -> Escenario -> Resultado", 24, True, WHITE)
tbl(sl, Inches(0.4), Inches(1.5), Inches(12.5), Inches(5.5),
    ["Requisito", "Funcionalidad", "Escenario Gherkin", "Resultado"],
    [["RF-01: Login valido", "Autenticacion", "TC_001_LoginExitoso", "PASS"],
     ["RF-02: Login invalido", "Error auth", "TC_002_LoginFallido", "PASS"],
     ["RF-03: Crear empleado", "PIM Registro", "TC_005_Agregar (3 filas DD)", "PASS"],
     ["RF-04: Buscar empleado", "PIM Busqueda", "TC_006_Buscar (3 filas DD)", "PASS"],
     ["RF-05: Editar empleado", "PIM Edicion", "TC_007_Editar (3 filas DD)", "PASS"],
     ["RF-06: Eliminar empleado", "PIM Eliminacion", "TC_008_Eliminar (3 filas DD)", "PASS"],
     ["RF-07: Navegacion", "Modulos", "TC_007-009_Nav (Leave/Admin/Info)", "PASS"],
     ["RF-08: Cerrar sesion", "Logout", "TC_003_Logout", "PASS"],
     ["RF-09: Validacion campos", "Formularios", "TC_010_ValidarCampos", "PASS"],
     ["RF-10: Assert titulo", "Verificacion", "TC_004_AssertTitulo", "FAIL (diseno)"],
     ["RF-11: Solicitar licencia", "Leave", "TC_014_LeaveApply (3 filas DD)", "PASS"]],
    cw=[Inches(2.5), Inches(2), Inches(4), Inches(1.5)])

# ========== SLIDE 5: RESULTADOS DE EJECUCION ==========
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE)
tb(sl, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8), "Resultados de Ejecucion", 32, True, WHITE)
tb(sl, Inches(0.8), Inches(1.5), Inches(5), Inches(0.5),
   "29 escenarios ejecutados | 96.55% aprobacion", 14, False, ACCENT)

bullets(sl, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5), [
    "Total escenarios: 29",
    "Aprobados: 28 (96.55%)",
    "Fallidos: 1 (3.45%) - TC_004 intencional",
    "Data-Driven: 18/18 PASS (62%)",
    "10 features ejecutadas",
    "Tiempo total: ~7 minutos",
    "177 screenshots en evidencias/",
    "6 archivos Excel testData",
], 16)

rect(sl, Inches(7), Inches(1.5), Inches(5.5), Inches(2.5), GREEN_BG)
tb(sl, Inches(7.3), Inches(1.7), Inches(5), Inches(0.4), "Fallo Intencional TC_004:", 14, True, GREEN)
tb(sl, Inches(7.3), Inches(2.1), Inches(5), Inches(1.6),
   "Escenario: Assert falla en verificacion de titulo\n"
   "Esperado: 'OrangeHRM OS 5.7'\n"
   "Real: 'OrangeHRM'\n"
   "Proposito: Demostrar hook @after_scenario\n"
   "captura screenshot automatico en fallos\n"
   "Evidencia en: evidencias/screenshot_*.png",
   12, False, BLUE)

rect(sl, Inches(7), Inches(4.3), Inches(5.5), Inches(2.5), RED_BG)
tb(sl, Inches(7.3), Inches(4.5), Inches(5), Inches(0.4), "Cobertura por Modulo:", 14, True, RED)
tb(sl, Inches(7.3), Inches(4.9), Inches(5), Inches(1.6),
   "Login: 3/4 PASS  |  Navegacion: 3/3 PASS\n"
   "Sesion: 1/1 PASS  |  PIM CRUD: 12/12 PASS\n"
   "Perfil: 3/3 PASS  |  Licencias: 3/3 PASS\n"
   "Data-Driven: 18/18 PASS  |  62% cobertura DD",
   12, False, BLUE)

# ========== SLIDE 6: METRICAS DE CALIDAD Y RENDIMIENTO ==========
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE)
tb(sl, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
   "Metricas de Calidad y Rendimiento", 28, True, WHITE)

tbl(sl, Inches(0.4), Inches(1.5), Inches(6), Inches(3.5),
    ["Metrica (PPT 3.3.1)", "Formula", "Resultado", "Fuente"],
    [["% Automatizable", "(Aut/Tot) x 100", "100% (29/29)", "reporte.html"],
     ["Progreso Autom.", "(Ejec/Planif) x 100", "100% (29/29)", "reporte.json"],
     ["Productiv. Diseno", "Esc / Horas", "7.25 esc/h", "Tiempo: 4 hrs"],
     ["Productiv. Ejec.", "Esc / T total", "0.067 esc/s", "reporte.json"],
     ["Tasa de Fallos", "(Fallos/Total) x 100", "3.45% (1/29)", "reporte.html"]],
    cw=[Inches(1.8), Inches(1.5), Inches(1.5), Inches(1.2)])

tbl(sl, Inches(6.8), Inches(1.5), Inches(6), Inches(3.5),
    ["Metrica Rendimiento", "Formula", "Resultado", "Fuente"],
    [["Tiempo prom./esc", "Suma dur / Total", "~15 seg", "reporte.json"],
     ["Feature mas lenta", "Max(dur feature)", "~40s (PIM)", "reporte.json"],
     ["Overhead setup", "Setup/Total", "~25%", "environment.py"],
     ["Cobertura DD", "DD / Total", "62% (18/29)", "Features"],
     ["Screenshots/esc", "Total SS / Esc", "6.1 (177/29)", "evidencias/"]],
    cw=[Inches(1.8), Inches(1.5), Inches(1.5), Inches(1.2)])

tb(sl, Inches(0.8), Inches(5.4), Inches(11), Inches(1.5),
   "Tasa de fallos: unico fallo = TC_004 intencional. Excluyendolo: 0%. Suite estable.\n"
   "Overhead setup: cada escenario abre su propio Chrome. Con paralelizacion (-j 4) se reduce.\n"
   "Cobertura DD 62%: escalable, permite agregar filas al Excel sin modificar codigo.",
   13, False, GRAY)

# ========== SLIDE 7: OPORTUNIDADES DE MEJORA ==========
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE)
tb(sl, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
   "Oportunidades de Mejora (A: Observacion -> B: Causa -> C: Accion)", 22, True, WHITE)

for i, (t, txt) in enumerate([
    ("O1: Menu en chino",
     "A) TC_007/008/009 fallaron\nB) Server locale cambia a chino\nC) MENU_TRANSLATIONS\n+ fallback bilingue implementado"),
    ("O2: Boton Save + Spinner",
     "A) Save interceptado por spinner\nB) Spinner OX tapa boton\nC) _wait_spinner_done()\n+ JS click implementado"),
    ("O3: Titulo incorrecto",
     "A) TC_004: titulo no coincide\nB) driver.title es 'OrangeHRM'\nC) Validar por URL en vez\nde titulo implementado"),
]):
    x = Inches(0.4 + i * 4.3)
    rect(sl, x, Inches(1.5), Inches(4), Inches(0.5), ACCENT, t, 12, True, WHITE, PP_ALIGN.CENTER)
    tb(sl, x, Inches(2.1), Inches(4), Inches(1.8), txt, 11, False, BLUE)

for i, (t, txt) in enumerate([
    ("O4: Tiempo ejecucion",
     "A) Suite ~7 minutos\nB) Escenarios secuenciales\n  driver nuevo c/u\nC) Propuesta: behave -j 4\n  (~2 min)"),
    ("O5: Excel no versionables",
     "A) 6 .xlsx editables manual\nB) Binarios, sin diff en git\nC) Propuesta: regenerar con\n  generar_excel.py pre-ejec"),
]):
    x = Inches(0.4 + i * 4.3)
    y = Inches(4.3)
    rect(sl, x, y, Inches(4), Inches(0.5), ACCENT, t, 12, True, WHITE, PP_ALIGN.CENTER)
    tb(sl, x, Inches(4.9), Inches(4), Inches(1.8), txt, 11, False, BLUE)

rect(sl, Inches(9), Inches(4.3), Inches(4), Inches(0.5), GREEN, "3 implementadas | 2 propuestas", 12, True, WHITE, PP_ALIGN.CENTER)

# ========== SLIDE 8: PROPUESTAS DE MEJORA ==========
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE)
tb(sl, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
   "Propuestas de Mejora (Dato -> Causa -> Accion -> Impacto)", 24, True, WHITE)

tbl(sl, Inches(0.4), Inches(1.5), Inches(12.5), Inches(3.5),
    ["Tipo", "Propuesta", "Dato", "Accion", "Impacto"],
    [["Producto", "MP-01: Mensajes\nerror diferenciados", "TC_002: msj generico\n'Invalid credentials'", "Separar casos:\nusuario vs pass", "UX + seguridad\ninformatica"],
     ["Producto", "MP-02: Selector\nidioma persistente", "INC-01/02: locale\ncambia a chino", "Cookie de idioma\npersistente", "UX internacional"],
     ["Proceso", "MP-03: Page\nObject Model", "Selectores en\n2+ archivos steps", "Refactor a clases:\nLoginPage, PimPage", "Mantenibilidad\ncentralizada"],
     ["Proceso", "MP-04: Ejecucion\nparalela", "Suite completa\n~7 minutos", "behave -j 4\n(4 workers)", "~2 min (-70%)"],
     ["Proceso", "MP-05: Logging\nestructurado", "print() comentados\nsin trazas", "Modulo logging\ncon timestamps", "Diagnostico\nrapido fallos"],
     ["Processo", "MP-06: CI/CD\nGitHub Actions", "Ejecucion\nmanual", "Workflow .github/\nen cada push", "Regresion\nautomatica"]],
    cw=[Inches(1.2), Inches(2), Inches(2.5), Inches(2.5), Inches(2)])

rect(sl, Inches(0.4), Inches(5.3), Inches(12.5), Inches(1.5), GREEN_BG)
tb(sl, Inches(0.6), Inches(5.4), Inches(12), Inches(1.3),
   "Criterio 3 (Producto OrangeHRM): MP-01, MP-02  |  Criterio 4 (Proceso Framework): MP-03, MP-04, MP-05, MP-06\n\n"
   "Cada propuesta nace de un dato real del reporte, explica la causa tecnica, propone un cambio concreto en el codigo, "
   "y cuantifica el impacto en el ciclo de vida del software.",
   12, False, BLUE)

# ========== SLIDE 9: STACK TECNOLOGICO ==========
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE)
tb(sl, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
   "Stack Tecnologico: Behave + Python + Selenium", 28, True, WHITE)

tbl(sl, Inches(0.4), Inches(1.5), Inches(12.5), Inches(3),
    ["Componente", "Version", "Proposito", "Configuracion"],
    [["Python", "3.14", "Lenguaje base", "Interprete: python.exe"],
     ["Behave", "1.2.6", "Framework BDD (Gherkin)", "environment.py + features/ + steps/"],
     ["Selenium", "4.28", "Autom. navegador Chrome", "WebDriverWait, By, EC"],
     ["openpyxl", "3.1+", "Lectura Excel DD", "ExcelUtils.get_cell_data()"],
     ["Chrome", "149", "Navegador bajo prueba", "ChromeDriver compatible"],
     ["PyCharm", "2026.1", "IDE desarrollo", "Run Config: behave features/"]],
    cw=[Inches(2), Inches(1.2), Inches(3.5), Inches(5)])

tb(sl, Inches(0.8), Inches(5), Inches(11), Inches(2),
   "Comandos principales:\n"
   "- python -m behave features/                 # Toda la suite\n"
   "- python -m behave features/login.feature    # Una feature\n"
   "- python -m behave --format json -o reporte.json   # Reporte JSON\n"
   "- set HEADLESS=0 && python -m behave features/     # Modo visual\n"
   "Gestion dependencias: pip install -r requirements.txt (equivalente a Maven)",
   13, False, BLUE)

# ========== SLIDE 10: CONCLUSIONES Y CIERRE ==========
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE)
tb(sl, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8), "Conclusiones y Proyecciones", 32, True, WHITE)

bullets(sl, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5), [
    "29/29 escenarios ejecutados",
    "96.55% aprobacion (fallo intencional TC_004)",
    "Suite bilingue: ingles + chino",
    "Data-Driven: 18 escenarios parametrizados (62%)",
    "177 screenshots de evidencia",
    "10 features, 6 archivos Excel",
    "Metricas trazables al reporte",
], 16)

rect(sl, Inches(7), Inches(1.6), Inches(5.5), Inches(2.5), DARK)
tb(sl, Inches(7.3), Inches(1.8), Inches(5), Inches(2.2),
   "Proyecciones:\n"
   "- GitHub Actions para CI/CD\n"
   "- Page Object Model completo\n"
   "- Allure Reports para tendencias\n"
   "- Ejecucion paralela behave -j 4\n"
   "- Expandir cobertura a mas modulos",
   13, False, WHITE)

rect(sl, Inches(7), Inches(4.5), Inches(5.5), Inches(2.5), GREEN_BG)
tb(sl, Inches(7.3), Inches(4.7), Inches(5), Inches(2),
   "5 oportunidades de mejora documentadas\n"
   "6 propuestas formales (2 producto + 4 proceso)\n"
   "Stack: Behave 1.2.6 + Selenium 4.28 + Python 3.14\n\n"
   "Lista para integracion en pipeline CI/CD",
   12, False, BLUE)

# ========== SLIDE 11: GRACIAS ==========
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, DARK)
rect(sl, Inches(0), Inches(3.2), Inches(13.333), Inches(0.08), ACCENT)
tb(sl, Inches(1), Inches(1.5), Inches(11), Inches(1.5), "Gracias", 48, True, WHITE, PP_ALIGN.CENTER)
tb(sl, Inches(1), Inches(3.8), Inches(11), Inches(2),
   "Luis Tasso  |  Seccion 802V\n"
   "API0101 - Automatizacion de Pruebas\n"
   "Prof. Manuel Soto",
   20, False, GRAY_TXT, PP_ALIGN.CENTER)

prs.save(OUTPUT)
print("PPT generado: " + OUTPUT)
