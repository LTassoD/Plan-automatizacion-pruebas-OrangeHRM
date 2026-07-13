"""Set Spanish proofing language on DOCX and PPTX."""
from docx import Document
from docx.oxml.ns import qn as wqn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.oxml.ns import qn as pqn
from lxml import etree
import copy

base = r'C:\Users\User\Desktop\Duoc\Semestre V\Automatizacion de Pruebas\AutomatizacionPruebasPython'

# === DOCX: set document default language ===
docx_path = f'{base}/INFORME_FINAL_EFT.docx'
doc = Document(docx_path)

# Set default language in styles.xml
styles = doc.styles.element
docDefaults = styles.find(wqn('w:docDefaults'))
if docDefaults is None:
    docDefaults = OxmlElement('w:docDefaults')
    styles.insert(0, docDefaults)

rPrDefault = docDefaults.find(wqn('w:rPrDefault'))
if rPrDefault is None:
    rPrDefault = OxmlElement('w:rPrDefault')
    docDefaults.append(rPrDefault)

rPr = rPrDefault.find(wqn('w:rPr'))
if rPr is None:
    rPr = OxmlElement('w:rPr')
    rPrDefault.append(rPr)

lang = rPr.find(wqn('w:lang'))
if lang is None:
    lang = OxmlElement('w:lang')
    rPr.append(lang)
lang.set(wqn('w:val'), 'es-CL')
lang.set(wqn('w:eastAsia'), 'es-CL')

# Also set lang on every run for compatibility
for para in doc.paragraphs:
    for run in para.runs:
        rPr = run._r.find(wqn('w:rPr'))
        if rPr is None:
            rPr = OxmlElement('w:rPr')
            run._r.insert(0, rPr)
        lang = rPr.find(wqn('w:lang'))
        if lang is None:
            lang = OxmlElement('w:lang')
            rPr.append(lang)
        lang.set(wqn('w:val'), 'es-CL')

# Same for tables
for table in doc.tables:
    for row in table.rows:
        for cell in row.cells:
            for para in cell.paragraphs:
                for run in para.runs:
                    rPr = run._r.find(wqn('w:rPr'))
                    if rPr is None:
                        rPr = OxmlElement('w:rPr')
                        run._r.insert(0, rPr)
                    lang = rPr.find(wqn('w:lang'))
                    if lang is None:
                        lang = OxmlElement('w:lang')
                        rPr.append(lang)
                    lang.set(wqn('w:val'), 'es-CL')

doc.save(docx_path)
print(f'DOCX: Saved with es-CL language')

# === PPTX: set language on every text run ===
pptx_path = f'{base}/PRESENTACION_FINAL_EFT.pptx'
prs = Presentation(pptx_path)

for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    rPr = run._r.find(pqn('a:rPr'))
                    if rPr is None:
                        rPr = OxmlElement('a:rPr')
                        run._r.insert(0, rPr)
                    lang = rPr.find(pqn('a:lang'))
                    if lang is None:
                        lang = rPr.makeelement(pqn('a:lang'), {})
                        rPr.append(lang)
                    lang.set('val', 'es-CL')
                # Also set default run properties for the paragraph
                pPr = para._p.find(pqn('a:pPr'))
                if pPr is None:
                    pPr = OxmlElement('a:pPr')
                    para._p.insert(0, pPr)
                defRPr = pPr.find(pqn('a:defRPr'))
                if defRPr is None:
                    defRPr = OxmlElement('a:defRPr')
                    pPr.append(defRPr)
                lang = defRPr.find(pqn('a:lang'))
                if lang is None:
                    lang = defRPr.makeelement(pqn('a:lang'), {})
                    defRPr.append(lang)
                lang.set('val', 'es-CL')

        # Handle tables
        try:
            if hasattr(shape, 'table') and shape.table is not None:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                rPr = run._r.find(pqn('a:rPr'))
                                if rPr is None:
                                    rPr = OxmlElement('a:rPr')
                                    run._r.insert(0, rPr)
                                lang = rPr.find(pqn('a:lang'))
                                if lang is None:
                                    lang = rPr.makeelement(pqn('a:lang'), {})
                                    rPr.append(lang)
                                lang.set('val', 'es-CL')
        except:
            pass

prs.save(pptx_path)
print(f'PPTX: Saved with es-CL language')
