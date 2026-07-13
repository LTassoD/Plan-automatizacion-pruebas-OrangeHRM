"""Recolor PPT with a modern dark theme."""
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt, Emu
from copy import deepcopy

# Table colors
TBL_HEADER_BG  = "1E293B"  # slate-800
TBL_HEADER_TXT = "FFFFFF"
TBL_ROW1_BG    = "2D3748"  # slate-750 (odd rows)
TBL_ROW2_BG    = "1E293B"  # slate-800 (even rows)
TBL_BODY_TXT   = "E2E8F0"  # slate-200

# === NEW COLOR PALETTE ===
C_DARK    = "111827"  # gray-900 (slides 1, 11 bg)
C_LIGHT   = "1E3A5F"  # azul oscuro sobrio (content bg)
C_ACCENT  = "818CF8"  # indigo-400 más claro (accent bar)
C_HEADING = "FFFFFF"  # white (titles on dark)
C_BODY    = "CBD5E1" # slate-300 (body text on dark)
C_WHITE   = "FFFFFF"

path = r'C:\Users\User\Desktop\Duoc\Semestre V\Automatizacion de Pruebas\AutomatizacionPruebasPython\PRESENTACION_FINAL_EFT.pptx'
prs = Presentation(path)

def set_bg(slide, color_hex):
    cSld = slide.background._element
    bg_elem = cSld.find(qn('p:bg'))
    if bg_elem is None:
        bg_elem = cSld.makeelement(qn('p:bg'), {})
        cSld.insert(0, bg_elem)
    bg_pr = bg_elem.find(qn('p:bgPr'))
    if bg_pr is None:
        bg_pr = bg_elem.makeelement(qn('p:bgPr'), {})
        bg_elem.append(bg_pr)
    for old in bg_pr.findall(qn('a:solidFill')):
        bg_pr.remove(old)
    for old in bg_pr.findall(qn('a:gradFill')):
        bg_pr.remove(old)
    solid = bg_pr.makeelement(qn('a:solidFill'), {})
    clr = solid.makeelement(qn('a:srgbClr'), {'val': color_hex})
    solid.append(clr)
    bg_pr.append(solid)

def set_fill(shape, color_hex):
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        spPr = sp.find(qn('a:spPr'))
    if spPr is None:
        return
    for old in spPr.findall(qn('a:solidFill')):
        spPr.remove(old)
    for old in spPr.findall(qn('a:gradFill')):
        spPr.remove(old)
    solid = spPr.makeelement(qn('a:solidFill'), {})
    clr = solid.makeelement(qn('a:srgbClr'), {'val': color_hex})
    solid.append(clr)
    spPr.insert(0, solid)

def set_text_color(shape, color_hex):
    """Recursively set color on all text runs."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            rPr = run._r.find(qn('a:rPr'))
            if rPr is None:
                rPr = run._r.makeelement(qn('a:rPr'), {})
                run._r.insert(0, rPr)
            for old in rPr.findall(qn('a:solidFill')):
                rPr.remove(old)
            solid = rPr.makeelement(qn('a:solidFill'), {})
            clr = solid.makeelement(qn('a:srgbClr'), {'val': color_hex})
            solid.append(clr)
            rPr.append(solid)

def set_table_color(cell, cell_bg, text_color):
    """Set cell background and text color in a table cell."""
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = cell._tc.makeelement(qn('a:tcPr'), {})
        cell._tc.insert(0, tcPr)
    for old in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(old)
    solid = tcPr.makeelement(qn('a:solidFill'), {})
    clr = solid.makeelement(qn('a:srgbClr'), {'val': cell_bg})
    solid.append(clr)
    tcPr.append(solid)

    # Set text color in all runs
    for para in cell.text_frame.paragraphs:
        pPr = para._p.find(qn('a:pPr'))
        if pPr is not None:
            defRPr = pPr.find(qn('a:defRPr'))
            if defRPr is not None:
                for old in defRPr.findall(qn('a:solidFill')):
                    defRPr.remove(old)
                solid2 = defRPr.makeelement(qn('a:solidFill'), {})
                clr2 = solid2.makeelement(qn('a:srgbClr'), {'val': text_color})
                solid2.append(clr2)
                defRPr.append(solid2)
        for run in para.runs:
            rPr = run._r.find(qn('a:rPr'))
            if rPr is None:
                rPr = run._r.makeelement(qn('a:rPr'), {})
                run._r.insert(0, rPr)
            for old in rPr.findall(qn('a:solidFill')):
                rPr.remove(old)
            solid2 = rPr.makeelement(qn('a:solidFill'), {})
            clr2 = solid2.makeelement(qn('a:srgbClr'), {'val': text_color})
            solid2.append(clr2)
            rPr.append(solid2)

def recolor_table(table):
    for row_idx, row in enumerate(table.rows):
        is_header = row_idx == 0
        bg = TBL_HEADER_BG if is_header else (TBL_ROW1_BG if row_idx % 2 == 1 else TBL_ROW2_BG)
        txt = TBL_HEADER_TXT if is_header else TBL_BODY_TXT
        for cell in row.cells:
            set_table_color(cell, bg, txt)

# === APPLY CHANGES ===
for idx, slide in enumerate(prs.slides):
    if idx == 0 or idx == len(prs.slides) - 1:
        # Slide 1 (Portada) and Slide 11 (Gracias): dark bg
        set_bg(slide, C_DARK)
    else:
        # Content slides: light bg
        set_bg(slide, C_LIGHT)

    # Recolor Rectangle 1 (accent bar on top of each slide)
    for shape in slide.shapes:
        if shape.name == 'Rectangle 1':
            set_fill(shape, C_ACCENT)

    # Recolor headings (first text shape with Heading 1 or large text)
    if idx == 0:
        # Cover: text stays white on dark bg
        for shape in slide.shapes:
            if shape.has_text_frame:
                set_text_color(shape, C_WHITE)
                # Add accent color to subtitle
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if 'API0101' in run.text or 'OrangeHRM' in run.text:
                            rPr = run._r.find(qn('a:rPr'))
                            if rPr is None:
                                rPr = run._r.makeelement(qn('a:rPr'), {})
                                run._r.insert(0, rPr)
                            for old in rPr.findall(qn('a:solidFill')):
                                rPr.remove(old)
                            solid = rPr.makeelement(qn('a:solidFill'), {})
                            clr = solid.makeelement(qn('a:srgbClr'), {'val': 'A5B4FC'})  # indigo-200
                            solid.append(clr)
                            rPr.append(solid)
    elif idx == len(prs.slides) - 1:
        # Gracias slide: white text
        for shape in slide.shapes:
            if shape.has_text_frame:
                set_text_color(shape, C_WHITE)
    else:
        # Content slides: white/light text on dark bg
        for shape in slide.shapes:
            if shape.shape_type == 19:  # TABLE
                recolor_table(shape.table)
            elif shape.has_text_frame:
                text = shape.text_frame.text.strip()
                # Titles get pure white, body text gets lighter gray
                is_title = False
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size and run.font.size >= Pt(24):
                            is_title = True
                            break
                    if is_title:
                        break
                if is_title or (shape.name.startswith('TextBox') and len(text) < 60 and not text.startswith('•')):
                    set_text_color(shape, C_WHITE)
                else:
                    set_text_color(shape, C_BODY)

out_path = path  # overwrite
prs.save(out_path)
print("PPT recolorizado exitosamente.")
