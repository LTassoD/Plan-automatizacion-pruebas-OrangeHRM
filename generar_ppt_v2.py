#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Genera PRESENTACION_EVALUACION_3.pptx segun guia BP_Presentaciones.pdf"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT = os.path.join(BASE_DIR, "PRESENTACION_EVALUACION_3.pptx")
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(0x2C, 0x3E, 0x50)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xE7, 0x4C, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x7F, 0x8C, 0x8D)
LIGHT = RGBColor(0xEC, 0xF0, 0xF1)
DARK = RGBColor(0x2C, 0x3E, 0x50)
ACCENT = RGBColor(0x29, 0x80, 0xB9)
GREEN_BG = RGBColor(0xE8, 0xF8, 0xF5)
RED_BG = RGBColor(0xFD, 0xED, 0xEC)
GRAY_TXT = RGBColor(0xBD, 0xC3, 0xC7)

def bg(sl, c=WHITE):
    f = sl.background.fill
    f.solid()
    f.fore_color.rgb = c

def rect(sl, l, t, w, h, c, text="", fs=14, b=False, fc=WHITE, a=PP_ALIGN.LEFT):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = c
    s.line.fill.background()
    if text:
        tf = s.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(fs)
        p.font.color.rgb = fc
        p.font.bold = b
        p.alignment = a
    return s

def tb(sl, l, t, w, h, text, fs=18, b=False, c=BLUE, a=PP_ALIGN.LEFT):
    s = sl.shapes.add_textbox(l, t, w, h)
    tf = s.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(fs)
    p.font.color.rgb = c
    p.font.bold = b
    p.alignment = a
    return s

def bullets(sl, l, t, w, h, items, fs=14, c=BLUE):
    s = sl.shapes.add_textbox(l, t, w, h)
    tf = s.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(fs)
        p.font.color.rgb = c
        p.space_after = Pt(6)
    return s

def tbl(sl, l, t, w, h, headers, rows, cw=None):
    ts = sl.shapes.add_table(len(rows)+1, len(headers), l, t, w, h)
    t = ts.table
    if cw:
        for i, w_ in enumerate(cw):
            t.columns[i].width = w_
    for i, h_ in enumerate(headers):
        c = t.cell(0, i)
        c.text = h_
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
        c.fill.solid()
        c.fill.fore_color.rgb = BLUE
    for ri, rd in enumerate(rows):
        for ci, v in enumerate(rd):
            c = t.cell(ri+1, ci)
            c.text = str(v)
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = BLUE
            if ri % 2 == 0:
                c.fill.solid()
                c.fill.fore_color.rgb = LIGHT
    return ts

# ========== SLIDE 1: PORTADA ==========
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, DARK)
rect(sl, Inches(0), Inches(2.5), Inches(13.333), Inches(0.08), ACCENT)
tb(sl, Inches(1), Inches(0.8), Inches(11), Inches(1), "EVALUACION PARCIAL 3", 40, True, WHITE, PP_ALIGN.CENTER)
tb(sl, Inches(1), Inches(1.6), Inches(11), Inches(0.8), "Analizando los Resultados Obtenidos", 28, False, GRAY_TXT, PP_ALIGN.CENTER)
tb(sl, Inches(1), Inches(3.2), Inches(11), Inches(3),
   "Automatizacion de Pruebas - OrangeHRM\n\n"
   "Estudiante: Luis Tasso  |  Seccion: 802V\n"
   "Profesor: Manuel Soto\n"
   "API0101 - Automatizacion de Pruebas",
   18, False, GRAY_TXT, PP_ALIGN.CENTER)

# ========== SLIDE 2: AGENDA ==========
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE)
tb(sl, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8), "Agenda", 32, True, WHITE)
bullets(sl, Inches(0.8), Inches(1.6), Inches(11), Inches(5), [
    "1. Trazabilidad Requisito -> Escenario -> Resultado (con reporte abierto)",
    "2. Resultados de Ejecucion - 29 escenarios, evidencias, fallo intencional TC_004",
    "3. Metricas de Calidad con origen verificable en el reporte",
    "4. Oportunidades de Mejora (A: Observacion -> B: Causa -> C: Accion)",
    "5. Propuestas al Producto (OrangeHRM) vs. Propuestas al Proceso (Framework)",
    "6. Lecciones Aprendidas y Conclusiones",
], 18)

# ========== SLIDE 3: TRAZABILIDAD ==========
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE)
tb(sl, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
   "Trazabilidad: Requisito -> Funcionalidad -> Escenario -> Resultado", 28, True, WHITE)
tbl(sl, Inches(0.4), Inches(1.5), Inches(12.5), Inches(5.5),
    ["Requisito", "Funcionalidad", "Escenario Gherkin", "Resultado"],
    [["RF-001: Login valido", "Autenticacion", "TC_001_LoginExitoso", "PASS"],
     ["RF-002: Login invalido", "Error auth", "TC_002_LoginFallido", "PASS"],
     ["RF-003: Campos vacios", "Validacion", "TC_010_ValidarCampos", "PASS"],
     ["RF-004: Logout", "Gestion sesion", "TC_003_Logout", "PASS"],
     ["RF-005: Nav. Leave", "Navegacion", "TC_007_NavegarLeave", "PASS"],
     ["RF-006: Nav. Admin", "Navegacion", "TC_008_NavegarAdmin", "PASS"],
     ["RF-007: Nav. My Info", "Navegacion", "TC_009_MyInfo", "PASS"],
     ["RF-008: CRUD empleados", "PIM", "TC_005/006/007/008", "PASS"],
     ["RF-009: Perfil My Info", "Perfil", "TC_013_MyInfo", "PASS"],
     ["RF-010: Licencias", "Leave", "TC_014_LeaveApply", "PASS"],
     ["RF-ABC: Assert titulo", "Verificacion", "TC_004_AssertTitulo", "FAIL (diseno)"]],
    cw=[Inches(2.2), Inches(2.2), Inches(3.5), Inches(1.5)])

# ========== SLIDE 4: RESULTADOS ==========
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE)
tb(sl, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8), "Resultados de Ejecucion", 32, True, WHITE)
tb(sl, Inches(0.8), Inches(1.6), Inches(5), Inches(0.6),
   "Reporte abierto: reporte.html - 29 escenarios ejecutados", 14, False, ACCENT)
bullets(sl, Inches(0.8), Inches(2.5), Inches(5.5), Inches(4), [
    "Total escenarios: 29",
    "Aprobados: 28 (96.55%)",
    "Fallidos: 1 (3.45%) - TC_004 intencional",
    "Features ejecutadas: 10 de 10",
    "Tiempo total: ~7 minutos",
    "133 screenshots en evidencias/",
    "6 archivos Excel Data-Driven",
], 16)
rect(sl, Inches(7), Inches(1.6), Inches(5.5), Inches(2.5), GREEN_BG)
tb(sl, Inches(7.3), Inches(1.8), Inches(5), Inches(0.4), "Fallo Intencional TC_004:", 14, True, GREEN)
tb(sl, Inches(7.3), Inches(2.2), Inches(5), Inches(1.6),
   "Escenario: Assert falla en verificacion de titulo\n"
   "Esperado: OrangeHRM OS 5.7 | Real: OrangeHRM\n"
   "Proposito: Demostrar hook @after_scenario\n"
   "captura screenshot automatico en fallos\n"
   "Evidencia: screenshot_20260616_204331.png",
   12, False, BLUE)
rect(sl, Inches(7), Inches(4.5), Inches(5.5), Inches(2.5), RED_BG)
tb(sl, Inches(7.3), Inches(4.7), Inches(5), Inches(0.4), "Cobertura por Modulo:", 14, True, RED)
tb(sl, Inches(7.3), Inches(5.1), Inches(5), Inches(1.6),
   "Login: 3/4 PASS  |  Navegacion: 3/3 PASS\n"
   "Sesion: 1/1 PASS  |  PIM CRUD: 12/12 PASS\n"
   "Perfil: 3/3 PASS  |  Licencias: 3/3 PASS\n"
   "Data-Driven: 18/18 PASS",
   12, False, BLUE)

# ========== SLIDE 5: METRICAS ==========
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE)
tb(sl, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
   "Metricas de Calidad (PPT 3.3.1) con Origen Verificable", 24, True, WHITE)
tbl(sl, Inches(0.4), Inches(1.5), Inches(12.5), Inches(4),
    ["Metrica", "Formula", "Resultado", "Fuente en el reporte"],
    [["% Automatizable", "(Automatizados / Totales) x 100", "100% (29/29)", "reporte.html - 29 ejecutados de 29"],
     ["Progreso Autom.", "(Ejecutados / Planificados) x 100", "100% (29/29)", "reporte.json - 10 features x 29"],
     ["Productiv. Diseno", "Escenarios / Horas diseno", "7.25 esc/h (29/4)", "Tiempo asignado: 4 hrs"],
     ["Productiv. Ejec.", "Escenarios / Tiempo total", "0.067 esc/s", "reporte.json - suma duraciones"],
     ["Tasa de Fallos", "(Fallos / Total) x 100", "3.45% (1/29)", "reporte.html - 28 PASS, 1 FAIL"]],
    cw=[Inches(2.5), Inches(3.5), Inches(2.5), Inches(4)])
tb(sl, Inches(0.8), Inches(5.8), Inches(11), Inches(1.2),
   "Todas las metricas apuntan a su fuente exacta en el reporte. La tasa de fallos corresponde unicamente a TC_004 (fallo intencional). Excluyendolo, la tasa es 0% - suite estable.",
   13, False, GRAY)

# ========== SLIDE 6: OPORTUNIDADES ==========
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE)
tb(sl, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
   "Oportunidades de Mejora (A: Observacion -> B: Causa -> C: Accion)", 22, True, WHITE)

for i, (t, txt) in enumerate([
    ("INC-01/02: Menu en chino",
     "A) TC_007/008/009 fallaron por menu en chino\nB) Server locale cambia, textos chinos\nC) MENU_TRANSLATIONS + fallback bilingue"),
    ("INC-04: Boton Save + Spinner",
     "A) Save interceptado por spinner OX\nB) Spinner tapa boton submit\nC) _wait_spinner_done() + JS click"),
    ("INC-03: Timeout PIM",
     "A) PIM carga lento (~40s/escenario)\nB) Suite secuencial, driver nuevo c/u\nC) Ejecucion paralela behave -j 4"),
]):
    x = Inches(0.4 + i * 4.3)
    rect(sl, x, Inches(1.5), Inches(4), Inches(0.5), ACCENT, t, 11, True, WHITE, PP_ALIGN.CENTER)
    tb(sl, x, Inches(2.1), Inches(4), Inches(1.5), txt, 11, False, BLUE)

for i, (t, txt) in enumerate([
    ("INC-05: Titulo incorrecto",
     "A) TC_004 falla: titulo esperado no coincide\nB) driver.title solo retorna OrangeHRM\nC) Validar por URL en vez de titulo"),
    ("Datos Excel no versionados",
     "A) 6 archivos Excel editables manualmente\nB) generar_excel.py existe, sin control\nC) Versionar + pre-commit hook"),
]):
    x = Inches(0.4 + i * 4.3)
    y = Inches(4.0)
    rect(sl, x, y, Inches(4), Inches(0.5), ACCENT, t, 11, True, WHITE, PP_ALIGN.CENTER)
    tb(sl, x, Inches(4.6), Inches(4), Inches(1.5), txt, 11, False, BLUE)

# ========== SLIDE 7: PROPUESTAS ==========
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE)
tb(sl, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
   "Propuestas de Mejora (Dato -> Causa -> Accion -> Impacto)", 24, True, WHITE)

tbl(sl, Inches(0.4), Inches(1.5), Inches(12.5), Inches(3.5),
    ["Tipo", "Propuesta", "Dato / Causa", "Accion", "Impacto"],
    [["Producto", "Mensajes error dif.", "TC_002: msj generico", "TC_002b con casos separados", "UX + seguridad"],
     ["Producto", "Localizacion persist.", "INC-01/02: locale", "Selector idioma en UI", "UX internacional"],
     ["Proceso", "Page Object Model", "Selectores en 2+ archivos", "Refactor a clases Page", "Mantenibilidad"],
     ["Proceso", "Ejecucion paralela", "Suite ~7 min, 15s/esc", "behave -j 4", "~2 min (-70%)"],
     ["Proceso", "Logging estructurado", "print() sin trazas", "Modulo logging estandar", "Diagnostico rapido"],
     ["Prod+Proc", "CI/CD GitHub Actions", "Ejecucion manual", "Workflow .github/", "Regresion automatica"]],
    cw=[Inches(1.2), Inches(2.2), Inches(2.8), Inches(3), Inches(2.2)])

rect(sl, Inches(0.4), Inches(5.3), Inches(12.5), Inches(1.5), GREEN_BG)
tb(sl, Inches(0.6), Inches(5.4), Inches(12), Inches(1.3),
   "Criterio 3 (Producto): MP-01, MP-02  |  Criterio 4 (Proceso): MP-03, MP-04, MP-05  |  Ambos: MP-06\n\n"
   "Cada propuesta nace de un dato real del reporte, explica la causa tecnica, propone un cambio concreto en el codigo, "
   "y cuantifica el impacto en el ciclo de vida del software (fase de mantenimiento, integracion continua o despliegue).",
   12, False, BLUE)

# ========== SLIDE 8: LECCIONES ==========
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE)
tb(sl, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8), "Lecciones Aprendidas", 32, True, WHITE)
bullets(sl, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5), [
    "Localizadores robustos = clave del exito en QA automation",
    "Data-Driven Testing: mas combinaciones, menos codigo",
    "Screenshots automaticos = trazabilidad total",
    "Internacionalizacion: un desafio real en entornos demo",
    "WebDriverWait > time.sleep() para estabilidad",
    "Suite 100% lista para integracion CI/CD",
], 16)
rect(sl, Inches(7), Inches(1.6), Inches(5.5), Inches(5), GREEN_BG)
tb(sl, Inches(7.3), Inches(1.8), Inches(5), Inches(0.4), "Reflexion personal:", 14, True, GREEN)
tb(sl, Inches(7.3), Inches(2.3), Inches(5), Inches(4),
   "Dificultades:\n"
   "- Locale cambiante del servidor demo\n"
   "- Tiempos carga variables (2s a 30s)\n"
   "- Encoding en decoradores behave\n\n"
   "Aprendizajes:\n"
   "- Fallo documentado > 10 pasos sin expl.\n"
   "- Trazabilidad = informe profesional\n"
   "- POM obligatorio al escalar\n\n"
   "Proyecciones:\n"
   "- GitHub Actions para CI\n"
   "- Migrar a Page Object classes\n"
   "- Allure Reports para tendencias",
   11, False, BLUE)

# ========== SLIDE 9: CONCLUSIONES ==========
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE)
tb(sl, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8), "Conclusiones", 32, True, WHITE)
bullets(sl, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5), [
    "29/29 escenarios ejecutados exitosamente",
    "96.55% de aprobacion (fallo intencional documentado)",
    "Suite robusta: soporta ingles y chino",
    "Data-Driven funcional con Excel (18/29 escenarios)",
    "133 evidencias visuales generadas",
    "Reporte HTML con metricas de calidad",
], 16)
rect(sl, Inches(7), Inches(1.6), Inches(5.5), Inches(5), DARK)
tb(sl, Inches(7.3), Inches(1.8), Inches(5), Inches(4.5),
   "La automatizacion de pruebas para\n"
   "OrangeHRM cumple con los objetivos\n"
   "de la evaluacion parcial 3:\n\n"
   "- 29 escenarios ejecutados\n"
   "- Evidencias (screenshots, reportes)\n"
   "- Oportunidades A+B+C\n"
   "- Metricas trazables al reporte\n"
   "- Propuestas Dato->Causa->Accion->Impacto\n\n"
   "Lista para CI/CD y regresion continua.",
   13, False, WHITE)

# ========== SLIDE 10: GRACIAS ==========
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, DARK)
rect(sl, Inches(0), Inches(3.2), Inches(13.333), Inches(0.08), ACCENT)
tb(sl, Inches(1), Inches(1.5), Inches(11), Inches(1.5), "Gracias", 48, True, WHITE, PP_ALIGN.CENTER)
tb(sl, Inches(1), Inches(3.8), Inches(11), Inches(2),
   "Luis Tasso  |  Seccion 802V\n"
   "API0101 - Automatizacion de Pruebas\n"
   "Prof. Manuel Soto",
   20, False, GRAY_TXT, PP_ALIGN.CENTER)

prs.save(OUTPUT)
print("PPT generado: " + OUTPUT)
